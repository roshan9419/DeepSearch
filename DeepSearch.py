import os
from time import sleep

try:
	import docx
	from pptx import Presentation
	from tkinter import *
	from tkinter import messagebox
	from tkinter import filedialog
except Exception as e:
	print('Requirement Not Satisfied')
	
def handleDOCXFile(file, search_key):
	doc = docx.Document(file)
	pr_no = 1
	location = 'Line : '
	flag = False
	for line in doc.paragraphs:
		if search_key in line.text.lower():
			flag = True
			# print('Found at Paragraph:', pr_no)
			# print(line.text)
			location += str(pr_no)+' '
		pr_no += 1

	if flag:
		return True, location
	return False, ''

def handlePPTXFile(file, search_key):
	prs = Presentation(file)
	slide_no = 1
	location = 'Slide: '
	flag = False
	for slide in prs.slides:
		for shapes in slide.shapes:
			if shapes.has_text_frame:
				# print(shapes.text)
				if(search_key in shapes.text.lower()):
					flag = True
					# print('Found at Slide', slide_no)
					location += str(slide_no) + ' '
					break
		slide_no += 1

	if flag:
		return True, location
	return False, ''

def handleOtherFiles(file, search_key):
	line_no=1
	location = 'Line : '
	flag = False
	with open(file, 'r') as f:
		for line in f:
			if search_key in line.lower():
				flag = True
				# print(file,'Line', line_no)
				location += str(line_no) + ' '
			line_no += 1

	if flag:
		return True, location
	return False, ''

def clearlist():
	resultList.delete(0, END)
	anotherList.delete(0, END)


i = 1
path = 'No Directory Choosen'
def getFiles():
	clearlist()
	global i, path
	# extensions = ['.pptx']
	extensions = ['.txt', '.py', '.docx', '.pptx', '.cpp']

	search_key = searchBox.get().lower()
	if path == 'No Directory Choosen':
		messagebox.showinfo("Error", "Choose the Directory")
		return
	if search_key == '':
		messagebox.showinfo("Invalid", "Enter in the Search Box")
		return
	flag = False
	os.chdir(path+'/')
	for file in os.listdir():
		f_name, f_ext = os.path.splitext(file)
		status = False
		location = 'None'
		if f_ext!='' and f_ext in extensions:
			if f_ext=='.docx':
				status, location = handleDOCXFile(file, search_key)
			elif f_ext=='.pptx':
				status, location = handlePPTXFile(file, search_key)
			else:
				status, location = handleOtherFiles(file, search_key)
		if(status):
			resultList.insert(i, file)
			anotherList.insert(i, location)
			flag = True
			i += 1
	if flag==False:
		messagebox.showinfo("Result", "Nothing Match")

	i = 1

def chooseDirectory():
	global path
	path = filedialog.askdirectory()
	if path == '':
		path = 'No Directory Choosen'
	fileLabel.configure(text=path)




root = Tk()
root.title('Deep Finder')
root.geometry('1000x500')


header = Label(root, text='Deep Finder', font=('Arial', 20), fg='red')
header.pack(pady=20)

searchBox = Entry(root, width=50, bd=2, font=('Arial', 18), justify='center', fg='gray') #searchBox.get()
searchBox.pack()

fileLabel = Label(root, text=path, font=('Arial', 10), fg='black')
fileLabel.pack()

browseBtn = Button(root, text='Browse', width=10, command=chooseDirectory)
browseBtn.pack(pady=10)

findBtn = Button(root, text='Find', width=10, command=getFiles, fg='white', bg='green')
findBtn.pack(pady=20)

frame = Frame(root)
frame.pack()

resultList = Listbox(frame, width=50)
resultList.grid(row=0, column=0)

anotherList = Listbox(frame, width=50)
anotherList.grid(row=0, column=1)


root.mainloop()


